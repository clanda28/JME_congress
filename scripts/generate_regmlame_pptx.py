from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

TITLE = "REGMLAME: Regularized Maximum Likelihood and Minimum Evolution"
SUBTITLE = "Speaker: <Your Name>  •  Affiliation  •  Date"

slides_plan = [
    {
        "title": "Title",
        "bullets": [],
        "notes": (
            "Introduce the motivation: structural ML bias and regularization as the remedy.\n"
            "Suggestions: Hero image of a phylogenetic tree; subtle background. No equations."
        ),
        "type": "title",
    },
    {
        "title": "Motivation: Why a new score?",
        "bullets": [
            "ML is powerful but structurally biased by branch-length–likelihood coupling",
            "Goal: stabilize inference via regularization",
        ],
        "notes": (
            "Emphasize bias is mathematical, not software-specific.\n"
            "Suggestions: Illustration contrasting ‘unbiased hope’ vs ‘structural bias’."
        ),
    },
    {
        "title": "Maximum Likelihood (ML) recap",
        "bullets": [
            "Optimize topology and branch lengths to maximize L",
            "Likelihood terms involve e^(tQ); branch lengths entangled with score",
        ],
        "notes": (
            "Set up the entanglement idea.\n"
            "Equations: P = exp(tQ)."
        ),
    },
    {
        "title": "Known ML failures",
        "bullets": [
            "Felsenstein: long-branch attraction contexts",
            "Kuhner and others: wrong topologies even with large data",
            "These are not rare edge cases",
        ],
        "notes": (
            "Cite classic simulation insights.\n"
            "Suggestions: Timeline/citation card graphic."
        ),
    },
    {
        "title": "Likelihood–branch length correlation",
        "bullets": [
            "Empirical r ≈ −0.9 to −0.95 across candidate trees",
            "Inflated branches can ‘buy’ higher likelihood",
        ],
        "notes": (
            "This correlation is the problem’s fingerprint.\n"
            "Graphics: Scatter plot: log-likelihood vs tree length; regression line."
        ),
    },
    {
        "title": "Overfitting symptoms in practice",
        "bullets": [
            "Over-resolved topologies",
            "Inflated external branches; noise/model violations absorbed",
        ],
        "notes": (
            "Audience has likely seen these patterns.\n"
            "Graphics: Side-by-side trees (balanced vs inflated tips)."
        ),
    },
    {
        "title": "Regularization analogy (statistics)",
        "bullets": [
            "Regression instability → L2 (ridge) penalty",
            "Penalize parameter magnitude to stabilize solutions",
        ],
        "notes": (
            "Bridge from regression to phylogenetics.\n"
            "Equations: Minimize ||y − Xβ||^2 + λ||β||^2."
        ),
    },
    {
        "title": "Defining REGMLAME",
        "bullets": [
            "Score = −log L + Λ Σ w_i b_i^2",
            "Λ tunes ML↔ME; Λ=0 → ML; Λ→∞ → Minimum Evolution",
        ],
        "notes": (
            "One knob to interpolate paradigms.\n"
            "Equations: Show the objective; annotate b_i, w_i, Λ."
        ),
    },
    {
        "title": "Unifying perspectives",
        "bullets": [
            "ML with penalty; ME with likelihood safeguard",
            "Bayesian MAP with short-branch prior",
            "Distance methods: same regularization logic",
        ],
        "notes": (
            "One framework, multiple interpretations.\n"
            "Graphics: Venn/bridge diagram linking ML, ME, MAP, Distance."
        ),
    },
    {
        "title": "Connection to distance methods",
        "bullets": [
            "Least-squares ≈ approximate likelihood",
            "Quadratic in branch lengths via distances mapping",
        ],
        "notes": (
            "Regularization naturally carries over.\n"
            "Equations: Classic LS objective over pairwise distances."
        ),
    },
    {
        "title": "Quadratic expansion and physics analogy",
        "bullets": [
            "Local harmonic approximation near ML optimum",
            "Keep diagonal Hessian terms for simplicity",
        ],
        "notes": (
            "Captures essential behavior with tractability.\n"
            "Equations: Taylor expansion: −log L(b) ≈ const + Σ a_i (b_i − b_i*)^2."
        ),
    },
    {
        "title": "Choosing Λ as a phase transition",
        "bullets": [
            "ML and ME as competing ‘phases’",
            "Vary Λ → transition; balance near critical region",
        ],
        "notes": (
            "Physics intuition guides calibration.\n"
            "Graphics: Two-well schematic; slider Λ moving balance."
        ),
    },
    {
        "title": "‘Specific heat’ diagnostic",
        "bullets": [
            "Plot derivative of regularized score vs Λ → peak at transition",
            "Select Λ near the peak",
        ],
        "notes": (
            "Robust, data-adaptive choice for Λ.\n"
            "Equations/graphics: C(Λ) ∝ d/dΛ [S(Λ)] with a peaked curve."
        ),
    },
    {
        "title": "Practical computation pipeline",
        "bullets": [
            "Get ML fits (topology, branch lengths, logL) from RAxML-NG/SSCPE",
            "Analyze logL vs tree length; infer branch-specific weights",
            "Solve for global Λ at transition point; compute REGMLAME",
        ],
        "notes": (
            "A posteriori scoring from standard outputs.\n"
            "Graphics: Flowchart of steps."
        ),
    },
    {
        "title": "Two variants of REGMLAME",
        "bullets": [
            "Non-rescaled: branch-specific, surgical penalties",
            "Rescaled: uniform after centering/rescaling; robust across datasets",
        ],
        "notes": (
            "Targeted vs robust penalty trade-off.\n"
            "Graphics: Side-by-side schematic of penalty profiles."
        ),
    },
    {
        "title": "Evaluation setup",
        "bullets": [
            "Protein families; compare ML, REGMLAME (non‑rescaled), REGMLAME (rescaled)",
            "Diagnostic: coefficient of variation across models/families",
        ],
        "notes": (
            "Stability as a key criterion.\n"
            "Graphics: Study design diagram."
        ),
    },
    {
        "title": "Variability results",
        "bullets": [
            "REGMLAME → lower coefficient of variation than ML",
            "More stable tree rankings",
        ],
        "notes": (
            "Regularization reduces volatility.\n"
            "Graphics: Bar chart of coefficients of variation (ML vs both variants)."
        ),
    },
    {
        "title": "Ranking differences",
        "bullets": [
            "ML favors its optimized models",
            "Non‑rescaled picks shorter/balanced branch trees (e.g., WAG/LG)",
            "Rescaled favors evenly distributed lengths (e.g., MF/DEWT)",
        ],
        "notes": (
            "REGMLAME reshuffles rankings, revealing ML bias.\n"
            "Graphics: Sankey/Alluvial diagram of best-tree selections by method."
        ),
    },
    {
        "title": "Limitations and next steps",
        "bullets": [
            "Currently a posteriori scoring",
            "Next: integrate penalty into branch-length optimization and topology search",
            "Start with distance methods (quadratic error)",
        ],
        "notes": (
            "Roadmap to full regularized inference.\n"
            "Graphics: Roadmap timeline."
        ),
    },
    {
        "title": "Conclusions",
        "bullets": [
            "ML bias is structural; regularization is the remedy",
            "REGMLAME unifies ML, ME, MAP, and distance methods",
            "Improves robustness and reproducibility",
        ],
        "notes": (
            "Close with the unifying message and call for adoption.\n"
            "Suggestions: Thank-you slide; optional QR code to repo/preprint."
        ),
    },
]

def add_slide_with_content(prs, title, bullets, notes, slide_index, total_slides, is_title=False):
    if is_title:
        layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = TITLE
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = SUBTITLE
    else:
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        if bullets:
            first = True
            for b in bullets:
                if first:
                    p = body.paragraphs[0]
                    p.text = b
                    p.level = 0
                    first = False
                else:
                    p = body.add_paragraph()
                    p.text = b
                    p.level = 0
    # Slide number footer (bottom-right)
    tf = slide.shapes.add_textbox(
        prs.slide_width - Inches(1.4),
        prs.slide_height - Inches(0.5),
        Inches(1.3),
        Inches(0.4),
    ).text_frame
    pnum = tf.paragraphs[0]
    pnum.text = f"{slide_index}/{total_slides}"
    pnum.font.size = Pt(10)
    pnum.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.clear()
    notes_tf.text = notes or ""

def build_presentation(output_path):
    prs = Presentation()
    # 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    total = len(slides_plan)
    for idx, spec in enumerate(slides_plan, start=1):
        add_slide_with_content(
            prs,
            title=spec.get("title", ""),
            bullets=spec.get("bullets", []),
            notes=spec.get("notes", ""),
            slide_index=idx,
            total_slides=total,
            is_title=(spec.get("type") == "title"),
        )

    prs.save(output_path)

if __name__ == "__main__":
    import os
    out_path = os.path.join("presentations", "REGMLAME.pptx")
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    build_presentation(out_path)